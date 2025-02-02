import streamlit as st
from langchain_ollama import OllamaLLM
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader, CSVLoader
from langchain.chains.question_answering import load_qa_chain
from langchain.chains.summarize import load_summarize_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from datetime import datetime
import json
import base64
import os
from tempfile import NamedTemporaryFile
from pptx import Presentation
import pyperclip
import time
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Streamlit app
st.set_page_config(page_title="Chatbot", layout="wide")

# Initialize session state variables
def initialize_session_state():
    if "all_chats" not in st.session_state:
        st.session_state.all_chats = [{"id": 0, "messages": [], "title": "New Chat", "file_vectorstore": None, "file_content": None}]
    if "current_chat_id" not in st.session_state:
        st.session_state.current_chat_id = 0
    if "regenerated_responses" not in st.session_state:
        st.session_state.regenerated_responses = {}
    if "feedback_data" not in st.session_state:
        st.session_state.feedback_data = []
    if "feedback_given" not in st.session_state:
        st.session_state.feedback_given = set()
    if "editing_message_idx" not in st.session_state:
        st.session_state.editing_message_idx = None
    if "show_settings" not in st.session_state:
        st.session_state.show_settings = False
    if "model_name" not in st.session_state:
        st.session_state.model_name = "deepseek-r1:1.5b"
    if "system_prompt" not in st.session_state:
        st.session_state.system_prompt = "You are a helpful assistant."
    if "next_chat_id" not in st.session_state:
        st.session_state.next_chat_id = 1
    if "show_delete_option" not in st.session_state:
        st.session_state.show_delete_option = None
    if "show_download_options" not in st.session_state:
        st.session_state.show_download_options = None

initialize_session_state()

# Get current chat
def get_current_chat():
    return next(chat for chat in st.session_state.all_chats if chat["id"] == st.session_state.current_chat_id)

current_chat = get_current_chat()

# Custom CSS with updated styles
st.markdown("""
    <style>
        [data-testid="stSidebar"][aria-expanded="true"] {
            width: 350px !important;
        }
        [data-testid="stSidebar"][aria-expanded="false"] {
            width: 80px !important;
            margin-left: 0px;
        }
        div[data-testid="stSidebarNav"] {
            display: none;
        }
        .download-button {
            display: inline-block;
            padding: 0.5em 1em;
            margin: 0.2em;
            border-radius: 4px;
            background-color: #FFFFFF;
            color: red;
            text-decoration: none;
            text-align: center;
        }
        .download-button:hover {
            background-color: #000000;
        }
    </style>
""", unsafe_allow_html=True)

# Initialize chatbot
def initialize_chatbot():
    return OllamaLLM(model=st.session_state.model_name)

llm = initialize_chatbot()
chat_prompt = ChatPromptTemplate.from_messages([
    ("system", st.session_state.system_prompt),
    ("human", "{input}")
])

# Initialize embeddings and text splitter
def initialize_embeddings():
    try:
        return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    except ImportError:
        st.error("The sentence-transformers package is not installed. Please install it with `pip install sentence-transformers`.")
        st.stop()

embeddings = initialize_embeddings()
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

# Helper functions
def show_temporary_message(container, message, is_error=False, duration=2):
    """Show a message temporarily in a container."""
    if is_error:
        container.error(message)
    else:
        container.success(message)
    time.sleep(duration)
    container.empty()

def process_pptx(file_path):
    """Extract text from a PowerPoint file."""
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def get_file_loader(file_extension, file_path):
    """Return the appropriate file loader based on the file extension."""
    if file_extension == ".pdf":
        return PyPDFLoader(file_path)
    elif file_extension == ".txt":
        return TextLoader(file_path)
    elif file_extension == ".docx":
        return Docx2txtLoader(file_path)
    elif file_extension == ".csv":
        return CSVLoader(file_path)
    elif file_extension == ".pptx":
        text = process_pptx(file_path)
        return [{"page_content": text, "metadata": {}}]  # Convert text to document format
    else:
        return None

def process_uploaded_file(file):
    """Process the uploaded file and create a vector store for it."""
    try:
        with NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_file_path = tmp_file.name

        file_extension = os.path.splitext(file.name)[1].lower()
        loader = get_file_loader(file_extension, tmp_file_path)

        if loader is None:
            st.error("Unsupported file type.")
            return

        if file_extension == ".pptx":
            documents = loader
        else:
            documents = loader.load()

        texts = text_splitter.split_documents(documents)
        vectorstore = FAISS.from_documents(texts, embeddings)
        current_chat["file_vectorstore"] = vectorstore
        current_chat["file_content"] = texts
        message_container = st.empty()
        show_temporary_message(message_container, "File processed successfully!")
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        message_container = st.empty()
        show_temporary_message(message_container, f"Error processing file: {e}", is_error=True)
    finally:
        os.unlink(tmp_file_path)

def answer_question_from_file(question):
    """Answer a question based on the uploaded file."""
    if "file_vectorstore" not in current_chat or current_chat["file_vectorstore"] is None:
        return None

    docs = current_chat["file_vectorstore"].similarity_search(question, k=3)
    chain = load_qa_chain(llm, chain_type="stuff")
    response = chain.run(input_documents=docs, question=question)
    return response

def summarize_file():
    """Generate a summary of the file content."""
    if "file_content" not in current_chat or current_chat["file_content"] is None:
        return "No file content available to summarize."

    chain = load_summarize_chain(llm, chain_type="map_reduce")
    summary = chain.run(current_chat["file_content"])
    return summary

def generate_response(prompt_text):
    """Generate a response using the chatbot."""
    formatted_prompt = chat_prompt.format(input=prompt_text)
    try:
        return llm.invoke(formatted_prompt)
    except Exception as e:
        logger.error(f"Error generating response: {e}")
        return "Sorry, I encountered an error. Please try again."

def create_new_chat():
    """Create a new chat with no file data."""
    new_chat_id = st.session_state.next_chat_id
    st.session_state.next_chat_id += 1
    new_chat = {
        "id": new_chat_id,
        "messages": [],
        "title": f"New Chat {new_chat_id}",
        "file_vectorstore": None,
        "file_content": None
    }
    st.session_state.all_chats.append(new_chat)
    st.session_state.current_chat_id = new_chat_id
    return new_chat

def get_file_status_message(chat):
    """Return a status message about file upload state."""
    if chat["file_content"] is not None:
        return "✅ File is loaded and ready for questions"
    return "No file uploaded in this chat"

# Sidebar content
with st.sidebar:
    settings_col1, settings_col2 = st.columns([1, 5])
    with settings_col1:
        if not st.session_state.show_settings:
            if st.button("⚙️", key="settings_button", help="Open Settings"):
                st.session_state.show_settings = True
                st.rerun()
        else:
            if st.button("✖️", key="close_settings", help="Close Settings"):
                st.session_state.show_settings = False
                st.rerun()
    
    with settings_col2:
        if st.button("+ New Chat", use_container_width=True):
            create_new_chat()
            st.rerun()

    if st.session_state.show_settings:
        st.markdown("### Settings")
        new_model = st.selectbox("Choose Model", ["deepseek-r1:1.5b", "other-model"], index=0)
        new_prompt = st.text_area("System Prompt", st.session_state.system_prompt)
        
        if st.button("Apply Settings", key="apply_settings"):
            if new_model != st.session_state.model_name:
                st.session_state.model_name = new_model
                llm = initialize_chatbot()
            
            if new_prompt != st.session_state.system_prompt:
                st.session_state.system_prompt = new_prompt
                chat_prompt = ChatPromptTemplate.from_messages([
                    ("system", new_prompt),
                    ("human", "{input}")
                ])
            st.success("Settings applied!")
            st.rerun()

        if st.button("Clear Current Chat", key="clear_chat"):
            current_chat["messages"] = []
            current_chat["file_vectorstore"] = None
            current_chat["file_content"] = None
            st.rerun()

    st.divider()
    
    # Display all chats
    st.markdown("### Recent Chats")
    for chat in reversed(st.session_state.all_chats):
        chat_title = chat["title"]
        is_current = chat["id"] == st.session_state.current_chat_id
        
        chat_col, menu_col = st.columns([4, 1])
        
        with chat_col:
            if st.button(
                f"{'→ ' if is_current else ''}{chat_title}",
                key=f"chat_{chat['id']}",
                use_container_width=True,
                type="secondary" if is_current else "primary"
            ):
                st.session_state.current_chat_id = chat["id"]
                st.rerun()
        
        with menu_col:
            if st.button("⋮", key=f"menu_{chat['id']}", help="Options"):
                st.session_state.show_delete_option = chat["id"]
                st.session_state.show_download_options = None
                st.rerun()
        
        # Show the delete and download options
        if st.session_state.show_delete_option == chat["id"]:
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("🗑️ Delete", key=f"delete_{chat['id']}"):
                    st.session_state.all_chats = [c for c in st.session_state.all_chats if c["id"] != chat["id"]]
                    if st.session_state.current_chat_id == chat["id"]:
                        if st.session_state.all_chats:
                            st.session_state.current_chat_id = st.session_state.all_chats[0]["id"]
                        else:
                            create_new_chat()
                    st.session_state.show_delete_option = None
                    st.rerun()
            
            with col2:
                if st.button("⬇️ Download", key=f"download_{chat['id']}"):
                    st.session_state.show_download_options = chat["id"]
                    st.rerun()
            
            # Show download format options
            if st.session_state.show_download_options == chat["id"]:
                download_col1, download_col2 = st.columns(2)
                
                with download_col1:
                    chat_text = "Chat Export (Text Format)\n" + "=" * 50 + "\n\n"
                    for msg in chat["messages"]:
                        role = "User" if msg["role"] == "human" else "Assistant"
                        chat_text += f"{role}:\n{msg['content']}\n\n"
                    txt_link = f'<a href="data:file/txt;base64,{base64.b64encode(chat_text.encode()).decode()}" download="chat_{chat["id"]}.txt" class="download-button">Download as TXT</a>'
                    st.markdown(txt_link, unsafe_allow_html=True)
                
                with download_col2:
                    chat_json = json.dumps({
                        "timestamp": datetime.now().isoformat(),
                        "messages": chat["messages"]
                    }, indent=2)
                    json_link = f'<a href="data:file/json;base64,{base64.b64encode(chat_json.encode()).decode()}" download="chat_{chat["id"]}.json" class="download-button">Download as JSON</a>'
                    st.markdown(json_link, unsafe_allow_html=True)

    # File uploader - specific to current chat
    st.divider()
    st.markdown("### File Upload")
    
    # Display current file status
    status_message = get_file_status_message(current_chat)
    st.markdown(f"**Current Status:** {status_message}")
    
    # Show file uploader
    uploaded_file = st.file_uploader("Upload a file", type=["pdf", "txt", "csv", "docx", "pptx"], key=f"file_uploader_{current_chat['id']}")
    
    if current_chat["file_content"] is not None:
        if st.button("🗑️ Clear Current File", key=f"clear_file_{current_chat['id']}"):
            current_chat["file_vectorstore"] = None
            current_chat["file_content"] = None
            message_container = st.empty()
            show_temporary_message(message_container, "File cleared from current chat")
            st.rerun()
    
    if uploaded_file is not None:
        process_uploaded_file(uploaded_file)
        message_container = st.empty()
        show_temporary_message(message_container, "File processed successfully!")

# Main chat interface
st.title("💬 Chatbot")

# Display chat messages
for idx, message in enumerate(current_chat["messages"]):
    role = "user" if message["role"] == "human" else "assistant"
    
    with st.chat_message(role):
        if idx == st.session_state.editing_message_idx:
            edit_cols = st.columns([4, 1])
            with edit_cols[0]:
                edited_text = st.text_input("Edit your question", value=message["content"], key=f"edit_input_{idx}")
            with edit_cols[1]:
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("✅", key=f"save_edit_{idx}"):
                        current_chat["messages"][idx]["content"] = edited_text
                        new_response = generate_response(edited_text)
                        if idx + 1 < len(current_chat["messages"]):
                            current_chat["messages"][idx + 1]["content"] = new_response
                        st.session_state.editing_message_idx = None
                        feedback_keys = [f"{idx+1}_positive", f"{idx+1}_negative"]
                        st.session_state.feedback_given = {k for k in st.session_state.feedback_given if k not in feedback_keys}
                        st.rerun()
                with col2:
                    if st.button("❌", key=f"cancel_edit_{idx}"):
                        st.session_state.editing_message_idx = None
                        st.rerun()
        else:
            cols = st.columns([10, 1])
            with cols[0]:
                st.markdown(message["content"])
            with cols[1]:
                if role == "user" and st.button("✏️", key=f"edit_{idx}", help="Edit question"):
                    st.session_state.editing_message_idx = idx
                    st.rerun()
        
        if role == "assistant":
            button_container = st.container()
            with button_container:
                cols = st.columns([1, 1, 1, 1, 5])  # Added one more column for copy button

                # Create empty containers for messages
                message_container = st.empty()
                
                with cols[0]:
                    if st.button("🔄", key=f"regenerate_{idx}", help="Regenerate response"):
                        if idx not in st.session_state.regenerated_responses:
                            st.session_state.regenerated_responses[idx] = []
                        st.session_state.regenerated_responses[idx].append(message["content"])
                        
                        new_response = generate_response(current_chat["messages"][idx - 1]["content"])
                        current_chat["messages"][idx]["content"] = new_response
                        feedback_keys = [f"{idx}_positive", f"{idx}_negative"]
                        st.session_state.feedback_given = {k for k in st.session_state.feedback_given if k not in feedback_keys}
                        st.rerun()
                
                with cols[1]:
                    feedback_key = f"{idx}_positive"
                    button_disabled = feedback_key in st.session_state.feedback_given
                    if st.button("👍", key=f"like_{idx}", disabled=button_disabled, help="Positive feedback"):
                        question = current_chat["messages"][idx - 1]["content"]
                        feedback_entry = {
                            "timestamp": datetime.now().isoformat(),
                            "question": question,
                            "response": message["content"],
                            "feedback": "positive",
                            "model": st.session_state.model_name,
                            "chat_id": st.session_state.current_chat_id
                        }
                        st.session_state.feedback_data.append(feedback_entry)
                        st.session_state.feedback_given.add(feedback_key)
                        show_temporary_message(message_container, "Feedback saved!")
                        st.rerun()
                
                with cols[2]:
                    feedback_key = f"{idx}_negative"
                    button_disabled = feedback_key in st.session_state.feedback_given
                    if st.button("👎", key=f"dislike_{idx}", disabled=button_disabled, help="Negative feedback"):
                        question = current_chat["messages"][idx - 1]["content"]
                        feedback_entry = {
                            "timestamp": datetime.now().isoformat(),
                            "question": question,
                            "response": message["content"],
                            "feedback": "negative",
                            "model": st.session_state.model_name,
                            "chat_id": st.session_state.current_chat_id
                        }
                        st.session_state.feedback_data.append(feedback_entry)
                        st.session_state.feedback_given.add(feedback_key)
                        show_temporary_message(message_container, "Feedback saved!")
                        st.rerun()
                
                with cols[3]:
                    if st.button("📋", key=f"copy_{idx}", help="Copy to clipboard"):
                        try:
                            pyperclip.copy(message["content"])
                            show_temporary_message(message_container, "Copied to clipboard!")
                        except Exception as e:
                            show_temporary_message(message_container, f"Failed to copy: {str(e)}")

# User input
if user_input := st.chat_input("Type your message...", key="main_input"):
    current_chat["messages"].append({"role": "human", "content": user_input})
    
    # Update chat title if it's the first message
    if len(current_chat["messages"]) == 1:
        current_chat["title"] = (user_input[:27] + "...") if len(user_input) > 30 else user_input
    
    # Handle specific commands
    if user_input.lower() == "summarize the file":
        if "file_content" in current_chat and current_chat["file_content"] is not None:
            summary = summarize_file()
            response = f"Summary of the file:\n\n{summary}"
        else:
            response = "No file has been uploaded to this chat. Please upload a file first."
    else:
        # Check if the question is related to the file
        if "file_vectorstore" in current_chat and current_chat["file_vectorstore"] is not None:
            file_response = answer_question_from_file(user_input)
            if file_response:
                response = file_response
            else:
                response = generate_response(user_input)
        else:
            response = generate_response(user_input)
    
    current_chat["messages"].append({"role": "assistant", "content": response})
    st.rerun()